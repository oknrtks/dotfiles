"""Shape操作ヘルパー関数

shape検索、座標取得、Placeholder長方形追加等の機能
"""

import logging
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


def find_shape_by_id(slide, shape_id: int):
    """shape_idからshapeを見つける

    Args:
        slide: スライドオブジェクト
        shape_id: 検索するshape_id

    Returns:
        見つかったshapeオブジェクト。見つからない場合はNone。
    """
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def get_shape_bounds(shape):
    """shapeの座標とサイズを取得

    Args:
        shape: シェイプオブジェクト

    Returns:
        (left, top, width, height) のタプル（EMU単位）
    """
    return (shape.left, shape.top, shape.width, shape.height)


def remove_shape_by_id(slide, shape_id: int) -> bool:
    """shape_idで指定したshapeを削除

    Args:
        slide: スライドオブジェクト
        shape_id: 削除するshape_id

    Returns:
        削除成功時はTrue、shapeが見つからない場合はFalse
    """
    shape = find_shape_by_id(slide, shape_id)
    if shape:
        slide.shapes._spTree.remove(shape.element)
        logging.info(f'Removed shape id={shape_id}')
        return True
    return False


def detect_text_style(slide, exclude_shape_ids: list = None) -> dict:
    """スライド内の主要なテキストshapeからスタイルを検出

    Args:
        slide: スライドオブジェクト
        exclude_shape_ids: 除外するshape_idのリスト

    Returns:
        検出されたスタイル情報の辞書 {'size': float, 'color': RGBColor, 'name': str}
    """
    if exclude_shape_ids is None:
        exclude_shape_ids = []

    styles = []

    for shape in slide.shapes:
        if shape.shape_id in exclude_shape_ids:
            continue

        if not shape.has_text_frame:
            continue

        text = shape.text.strip()
        if not text or text.startswith('【'):
            continue

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and run.font.size.pt >= 10:
                    style = {
                        'size': run.font.size.pt,
                        'color': None,
                        'name': run.font.name,
                    }

                    try:
                        if run.font.color.rgb:
                            style['color'] = run.font.color.rgb
                    except:
                        pass

                    styles.append(style)
                    break
            break

    if not styles:
        return {'size': 11, 'color': RGBColor(0x55, 0x55, 0x55), 'name': None}

    sizes = [s['size'] for s in styles]
    most_common_size = max(set(sizes), key=sizes.count)

    for style in styles:
        if style['size'] == most_common_size:
            if style['color'] is None:
                style['color'] = RGBColor(0x55, 0x55, 0x55)
            return style

    return styles[0]


def add_placeholder_rect(slide, left, top, width, height, text: str, style: dict = None):
    """Placeholder長方形を追加（図領域代替用、角丸なし）

    視認性を確保するため、薄いグレー背景とグレー枠線を設定。
    周囲のテキストスタイルを適用可能。

    Args:
        slide: スライドオブジェクト
        left: 左位置（EMU単位）
        top: 上位置（EMU単位）
        width: 幅（EMU単位）
        height: 高さ（EMU単位）
        text: Placeholderテキスト
        style: スタイル辞書 {'size': float, 'color': RGBColor, 'name': str}

    Returns:
        作成されたshapeオブジェクト
    """
    if style is None:
        style = {'size': 11, 'color': RGBColor(0x55, 0x55, 0x55), 'name': None}

    # 純粋な長方形を作成（角丸なし）
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )

    # 角丸を無効化（エラーハンドリング付き）
    try:
        if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
            shape.adjustments[0] = 0
    except:
        pass

    # 薄いグレー背景（視認性確保）
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

    # グレー枠線
    shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
    shape.line.width = Pt(1)

    # テキスト設定（スタイルを適用）
    if shape.has_text_frame:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        run = p.add_run()
        run.text = text

        # スタイルを適用
        run.font.size = Pt(style['size'])

        if style.get('color'):
            run.font.color.rgb = style['color']
        else:
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        if style.get('name'):
            run.font.name = style['name']

    logging.info(f'Added placeholder rect: size={style["size"]}pt')
    return shape
