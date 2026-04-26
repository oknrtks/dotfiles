from typing import Optional
from pptx.table import Table
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
import logging


def extract_table_structure(table_shape) -> dict:
    """表の構造情報を抽出する"""
    table = table_shape.table
    rows = len(table.rows)
    cols = len(table.columns)

    header_row_exists = False
    header_bg_color = None

    if rows > 0:
        first_row_cells = [table.cell(0, col_idx) for col_idx in range(cols)]
        bold_count = 0
        bg_colors = set()

        for cell in first_row_cells:
            if cell.text_frame.paragraphs:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.bold:
                            bold_count += 1
                            break

            if cell.fill.type == 1:  # SOLID fill
                try:
                    color = cell.fill.fore_color.rgb
                    bg_colors.add(str(color))
                except:
                    pass

        if bold_count >= cols / 2 or len(bg_colors) > 0:
            header_row_exists = True
            if bg_colors:
                header_bg_color = list(bg_colors)[0]

    cell_formats = []
    for row_idx in range(rows):
        row_formats = []
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            cell_format = {
                'font_size': None,
                'bold': False,
                'color': None,
                'bg_color': None
            }

            if cell.text_frame.paragraphs:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            cell_format['font_size'] = run.font.size.pt
                        if run.font.bold:
                            cell_format['bold'] = True
                        try:
                            if run.font.color.rgb:
                                cell_format['color'] = str(run.font.color.rgb)
                        except:
                            pass
                        break
                    break

            if cell.fill.type == 1:
                try:
                    cell_format['bg_color'] = str(cell.fill.fore_color.rgb)
                except:
                    pass

            row_formats.append(cell_format)
        cell_formats.append(row_formats)

    return {
        'rows': rows,
        'cols': cols,
        'header_row_exists': header_row_exists,
        'header_bg_color': header_bg_color,
        'cell_formats': cell_formats
    }


def placeholderize_table_cell(cell, placeholder_text: str, font_size_pt: int = None) -> None:
    """表セルのテキストをPlaceholder化する（書式保持、フォントサイズ指定可能）

    Args:
        cell: 表のセルオブジェクト
        placeholder_text: 設定するPlaceholderテキスト
        font_size_pt: フォントサイズ（pt）。Noneの場合は元のサイズを保持
    """
    text_frame = cell.text_frame

    if not text_frame.paragraphs:
        return

    first_para = text_frame.paragraphs[0]

    if first_para.runs:
        first_run = first_para.runs[0]

        for run in list(first_para.runs[1:]):
            run._element.getparent().remove(run._element)

        first_run.text = placeholder_text

        # フォントサイズを指定された値に変更
        if font_size_pt:
            first_run.font.size = Pt(font_size_pt)
    else:
        run = first_para.add_run()
        run.text = placeholder_text
        if font_size_pt:
            run.font.size = Pt(font_size_pt)

    for para in list(text_frame.paragraphs[1:]):
        para._element.getparent().remove(para._element)


def convert_table_to_placeholder(slide, table_shape_id: int) -> bool:
    """表のセル内容をPlaceholder化する（構造は保持、詳細情報付き）"""
    from .shapes import find_shape_by_id

    table_shape = find_shape_by_id(slide, table_shape_id)
    if not table_shape or not table_shape.has_table:
        return False

    structure = extract_table_structure(table_shape)
    table = table_shape.table

    logging.info(f'Table structure: {structure["rows"]}x{structure["cols"]}, '
                 f'header={structure["header_row_exists"]}')

    # ヘッダー行の処理
    if structure['header_row_exists'] and structure['rows'] > 0:
        for col_idx in range(structure['cols']):
            cell = table.cell(0, col_idx)
            placeholderize_table_cell(cell, f'ヘッダ{col_idx + 1}')
        start_row = 1
    else:
        start_row = 0

    # データ行の処理 - 最初のセルに詳細情報を必ず含める
    for row_idx in range(start_row, structure['rows']):
        for col_idx in range(structure['cols']):
            cell = table.cell(row_idx, col_idx)

            # データ行の最初のセルに詳細情報を含める
            if col_idx == 0 and row_idx == start_row:
                # デフォルトの詳細情報
                placeholder_text = f"""この表にデータを入力する
・行数: {structure['rows']}行 / 列数: {structure['cols']}列
・ヘッダー行: {'あり' if structure['header_row_exists'] else 'なし'}
・データ行書式: 文字色: #000000 / サイズ: 10.5pt
・セル背景: 交互色または単色
・文体: ですます調 / 体言止め / 簡潔表記
・特記: 元スライドでは表が使用されていた"""
                placeholderize_table_cell(cell, placeholder_text, font_size_pt=9)
                logging.info(f'Added detailed placeholder to cell [{row_idx},{col_idx}]')
            else:
                placeholderize_table_cell(cell, '【データ】')

    return True


def convert_table_to_placeholder_with_details(slide, table_shape_id: int, table_info: dict) -> bool:
    """表のセル内容をPlaceholder化（データ行の最初のセルに詳細情報を含める）

    Args:
        slide: スライドオブジェクト
        table_shape_id: 表のshape_id
        table_info: 表の情報を含む辞書
            - color: データ行の文字色（例: '#000000'）
            - size: データ行のフォントサイズ（例: '10.5'）
            - background: セル背景（例: '交互色（薄青・薄紫）'）
            - style: 文体（例: '箇条書き / ですます調'）
            - note: 特記事項（例: '元スライドでは進捗表が使用されていた'）

    Returns:
        成功時True、失敗時False
    """
    from .shapes import find_shape_by_id

    table_shape = find_shape_by_id(slide, table_shape_id)
    if not table_shape or not table_shape.has_table:
        return False

    structure = extract_table_structure(table_shape)
    table = table_shape.table

    logging.info(f'Table structure: {structure["rows"]}x{structure["cols"]}, '
                 f'header={structure["header_row_exists"]}')

    # ヘッダー行の処理
    if structure['header_row_exists'] and structure['rows'] > 0:
        for col_idx in range(structure['cols']):
            cell = table.cell(0, col_idx)
            placeholderize_table_cell(cell, f'ヘッダ{col_idx + 1}')
        start_row = 1
    else:
        start_row = 0

    # データ行の処理
    for row_idx in range(start_row, structure['rows']):
        for col_idx in range(structure['cols']):
            cell = table.cell(row_idx, col_idx)

            # データ行の最初のセルに詳細情報を含める
            if col_idx == 0 and row_idx == start_row:
                placeholder_text = f"""この表にデータを入力する
・行数: {structure['rows']}行 / 列数: {structure['cols']}列
・ヘッダー行: {'あり' if structure['header_row_exists'] else 'なし'}
・データ行書式: 文字色: {table_info.get('color', '#000000')} / サイズ: {table_info.get('size', '10.5')}pt
・セル背景: {table_info.get('background', '交互色または単色')}
・文体: {table_info.get('style', 'ですます調 / 体言止め / 簡潔表記')}
・特記: {table_info.get('note', '元スライドではデータ表が使用されていた')}"""
                # 長いテキストなので小さいフォントで
                placeholderize_table_cell(cell, placeholder_text, font_size_pt=9)
                logging.info(f'Added detailed placeholder to cell [{row_idx},{col_idx}]')
            else:
                placeholderize_table_cell(cell, '【データ】')

    return True
