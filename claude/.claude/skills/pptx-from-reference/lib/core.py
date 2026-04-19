"""PPTXテンプレート作成コア関数

書式保持テキスト置換、スライド操作等の基本機能
"""

import logging
from pptx import Presentation

R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def set_text_preserve_style(shape, text: str) -> None:
    """シェイプの書式を保持しつつテキストを置換

    text_frame.clear()は絶対に使用せず、先頭runのみを置換する。
    これにより、フォント・サイズ・色などの書式情報が保持される。

    Args:
        shape: テキストを設定するシェイプオブジェクト
        text: 設定するテキスト
    """
    tf = shape.text_frame
    # 先頭段落・先頭ランのみ残してテキストを設定
    first_para = tf.paragraphs[0]

    # 2段落目以降を削除
    for para in list(tf.paragraphs[1:]):
        tf._txBody.remove(para._p)

    # 先頭runを再利用
    runs = first_para.runs
    if runs:
        runs[0].text = text
        # 2run目以降を削除
        for run in list(runs[1:]):
            first_para._p.remove(run._r)
    else:
        first_para.add_run().text = text

    logging.info(f'Set text on shape id={shape.shape_id}: {text!r}')


def remove_slide(prs: Presentation, index: int) -> None:
    """指定インデックスのスライドを削除

    注意: 複数スライドを削除する場合は、後ろのインデックスから順に実行すること。

    Args:
        prs: プレゼンテーションオブジェクト
        index: 削除するスライドのインデックス（0-indexed）
    """
    sll = prs.slides._sldIdLst
    sld_id = sll[index]
    rId = sld_id.attrib.get(f'{{{R_NS}}}id')
    sll.remove(sld_id)
    prs.part.drop_rel(rId)
    logging.info(f'Removed slide index={index} rId={rId}')


def clone_slide_to_end(prs: Presentation, index: int) -> None:
    """指定インデックスのスライドを複製して末尾に追加

    Args:
        prs: プレゼンテーションオブジェクト
        index: 複製するスライドのインデックス（0-indexed）
    """
    import copy

    src = prs.slides[index]
    new_slide = prs.slides.add_slide(src.slide_layout)

    # シェイプツリーをコピー
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in src.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    # 画像リレーションをコピー
    for rel in src.part.rels.values():
        if 'image' in rel.reltype:
            new_slide.part.rels._rels[rel.rId] = rel

    logging.info(f'Cloned slide index={index} → appended as slide {len(prs.slides)}')


def move_last_slide_to_front(prs: Presentation) -> None:
    """末尾スライドを先頭に移動

    Args:
        prs: プレゼンテーションオブジェクト
    """
    sll = prs.slides._sldIdLst
    last = sll[-1]
    sll.remove(last)
    sll.insert(0, last)
    logging.info('Moved last slide to front')
