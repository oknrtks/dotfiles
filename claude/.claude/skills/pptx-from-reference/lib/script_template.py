#!/usr/bin/env python3
"""テンプレート生成スクリプト - 自動生成用テンプレート

このファイルはテンプレートであり、LLMが各プロジェクト用にカスタマイズする。
"""

import sys
from pathlib import Path
from pptx import Presentation
from datetime import datetime
import logging

# ライブラリパスを設定
sys.path.insert(0, str(Path.home() / '.claude' / 'skills' / 'pptx-from-reference'))
from lib.core import set_text_preserve_style, remove_slide
from lib.shapes import find_shape_by_id, add_placeholder_rect, remove_shape_by_id

# ========================================
# ログ設定
# ========================================
logging.basicConfig(
    filename=f'_work/build_{datetime.now().strftime("%Y%m%d_%H%M%S")}.log',
    level=logging.DEBUG,
    format='%(asctime)s %(levelname)s %(message)s',
)

# ========================================
# ファイル固有変数の定義
# ========================================
# 以下の変数は layout_memo.md の分析結果に基づいて設定する

INPUT = Path('INPUT.pptx')  # TODO: 入力ファイル名に置換
OUTPUT = Path('template_INPUT.pptx')  # TODO: 出力ファイル名に置換

# 保持するスライドインデックス（0-indexed）
KEEP_INDICES = {0, 1, 3, 7}  # TODO: layout_memo.mdに基づいて設定

# 各スライドのshape_idマッピング
# TODO: layout_memo.md の値をここに転記
SLIDE_CONFIG = {
    0: {  # スライド1
        'title_id': None,  # TODO: shape_idを設定
        'subtitle_id': None,  # TODO: 必要に応じて追加
    },
    1: {  # スライド2
        'title_id': None,  # TODO
        'lead_id': None,  # TODO
        'body_ids': [],  # TODO: 複数ある場合はリスト
    },
    # ... 他スライドも同様
}

# Placeholderテキストの定義
# placeholder_templates.md からコピーして、#XXXXXX等を置換
PLACEHOLDER_TEXTS = {
    'title': '【タイトル】スライドタイトルをここに記載',
    'subtitle': '【表紙サブタイトル】ここにサブタイトル・部署名等を記載',
    'lead': '【リード文】このスライドの要点・キーメッセージを1文で記載',
    'body': '【本文】詳細説明・根拠・データなどを記載\n'
            '・文字色: 元の書式を継承 / サイズ: XXpt\n'
            '・強調方法: Bold / 文体: ですます調\n'
            '・構成: 1ペイン / 図: なし / 表: なし',
    'figure_area': '【図領域】この箇所には概念図・グラフ・イメージ図を配置する\n'
                  '（元スライドで XXX が使用されていた）',
}

# ========================================
# ヘルパー関数（プロジェクト固有）
# ========================================
def process_slide(slide, config):
    """スライドの処理（テキスト置換・Shape操作）

    Args:
        slide: スライドオブジェクト
        config: SLIDE_CONFIGの該当スライドの設定辞書
    """
    # タイトル設定
    if 'title_id' in config and config['title_id']:
        shape = find_shape_by_id(slide, config['title_id'])
        if shape:
            set_text_preserve_style(shape, PLACEHOLDER_TEXTS['title'])

    # サブタイトル設定
    if 'subtitle_id' in config and config['subtitle_id']:
        shape = find_shape_by_id(slide, config['subtitle_id'])
        if shape:
            set_text_preserve_style(shape, PLACEHOLDER_TEXTS['subtitle'])

    # リード文設定
    if 'lead_id' in config and config['lead_id']:
        shape = find_shape_by_id(slide, config['lead_id'])
        if shape:
            set_text_preserve_style(shape, PLACEHOLDER_TEXTS['lead'])

    # 本文設定（複数ある場合）
    if 'body_ids' in config:
        for body_id in config['body_ids']:
            shape = find_shape_by_id(slide, body_id)
            if shape:
                set_text_preserve_style(shape, PLACEHOLDER_TEXTS['body'])

    # 図領域の処理（必要な場合）
    # if 'figure_shape_id' in config:
    #     img_shape = find_shape_by_id(slide, config['figure_shape_id'])
    #     if img_shape:
    #         coords = (img_shape.left, img_shape.top, img_shape.width, img_shape.height)
    #         remove_shape_by_id(slide, config['figure_shape_id'])
    #         add_placeholder_rect(slide, *coords, PLACEHOLDER_TEXTS['figure_area'])


# ========================================
# メイン処理
# ========================================
def main():
    logging.info('テンプレート生成開始')

    prs = Presentation(OUTPUT)
    logging.info(f'Opened: {len(prs.slides)} slides')

    # 不要スライドを削除（後ろから）
    for i in range(len(prs.slides) - 1, -1, -1):
        if i not in KEEP_INDICES:
            remove_slide(prs, i)
            logging.info(f'Removed slide index={i}')

    logging.info(f'After deletion: {len(prs.slides)} slides')

    # 各スライドの処理
    for i in range(len(prs.slides)):
        slide = prs.slides[i]
        config = SLIDE_CONFIG.get(i, {})
        logging.info(f'Processing slide {i + 1}')
        process_slide(slide, config)

    prs.save(OUTPUT)
    logging.info(f'Saved: {OUTPUT}')
    print(f'Done: {OUTPUT}')


if __name__ == '__main__':
    main()
