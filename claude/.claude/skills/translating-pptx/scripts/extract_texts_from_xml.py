#!/usr/bin/env python
"""
PPTXからテキストを抽出するスクリプト（汎用）

任意のPPTXファイルから全スライドのテキストを抽出し、
スライドごとにJSONファイルとして保存します。

使用方法:
    python extract_texts_from_xml.py <input.pptx> <output_dir>

出力:
    slide1_texts.json, slide2_texts.json, ...
"""

from pptx import Presentation
import json
import sys
from pathlib import Path


def has_japanese(text: str) -> bool:
    """
    テキストに日本語が含まれるか判定

    Args:
        text: 判定するテキスト

    Returns:
        日本語が含まれる場合はTrue
    """
    for char in text:
        if '\u3040' <= char <= '\u309F':  # ひらがな
            return True
        if '\u30A0' <= char <= '\u30FF':  # カタカナ
            return True
        if '\u4E00' <= char <= '\u9FFF':  # 漢字
            return True
    return False


def extract_slide_texts(slide, slide_idx: int) -> dict:
    """
    1つのスライドからテキストを抽出（テーブル対応版）

    Args:
        slide: pptx.slide.Slideオブジェクト
        slide_idx: スライド番号（1から開始）

    Returns:
        スライドのテキストデータ（辞書形式）
    """
    texts = []

    for shape_idx, shape in enumerate(slide.shapes):
        # 通常のテキストフレームを処理
        if hasattr(shape, "text_frame"):
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(para.runs):
                    if run.text.strip():
                        texts.append({
                            'shape_idx': shape_idx,
                            'para_idx': para_idx,
                            'run_idx': run_idx,
                            'text': run.text,
                            'has_japanese': has_japanese(run.text)
                        })

        # テーブルを処理
        if shape.has_table:
            table = shape.table
            for row_idx in range(len(table.rows)):
                row = table.rows[row_idx]
                for col_idx in range(len(row.cells)):
                    cell = row.cells[col_idx]
                    if cell.text_frame:
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                if run.text.strip():
                                    texts.append({
                                        'shape_idx': shape_idx,
                                        'table_row': row_idx,
                                        'table_col': col_idx,
                                        'para_idx': para_idx,
                                        'run_idx': run_idx,
                                        'text': run.text,
                                        'has_japanese': has_japanese(run.text),
                                        'is_table_text': True
                                    })

    return {
        'slide_number': slide_idx,
        'total_texts': len(texts),
        'total_japanese_texts': sum(1 for t in texts if t['has_japanese']),
        'texts': texts
    }


def extract_all_slide_texts(input_pptx: str, output_dir: str):
    """
    任意のPPTXファイルから全スライドのテキストを抽出

    Args:
        input_pptx: PPTXファイルパス
        output_dir: 出力ディレクトリ

    出力:
        slide1_texts.json, slide2_texts.json, ...
    """
    input_path = Path(input_pptx).resolve()
    output_path = Path(output_dir).resolve()
    output_path.mkdir(parents=True, exist_ok=True)

    print(f"Extracting texts from {input_path.name}...")

    # PPTXを読み込み
    prs = Presentation(input_pptx)

    print(f"Total slides: {len(prs.slides)}\n")

    # 各スライドのテキストを抽出
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = extract_slide_texts(slide, slide_idx)

        # JSONファイルとして保存
        output_file = output_path / f"slide{slide_idx}_texts.json"
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(slide_data, f, ensure_ascii=False, indent=2)

        print(f"Slide {slide_idx}:")
        print(f"  Total texts: {slide_data['total_texts']}")
        print(f"  Japanese texts: {slide_data['total_japanese_texts']}")
        print(f"  Saved to: {output_file.name}")
        print()

    print(f"✓ Successfully extracted texts from {len(prs.slides)} slides")
    print(f"Output directory: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_texts_from_xml.py <input.pptx> <output_dir>")
        print("\nExample:")
        print("  python extract_texts_from_xml.py sample.pptx extracted/")
        sys.exit(1)

    input_pptx = sys.argv[1]
    output_dir = sys.argv[2]

    if not Path(input_pptx).exists():
        print(f"Error: File not found: {input_pptx}")
        sys.exit(1)

    extract_all_slide_texts(input_pptx, output_dir)
