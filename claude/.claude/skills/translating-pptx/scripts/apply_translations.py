#!/usr/bin/env python
"""
対訳をPPTXに適用するスクリプト（汎用）

AIが作成した対訳JSONファイルを使用して、
任意のPPTXファイルに翻訳を適用し、色付け（RGB:128,0,128）を行います。

特徴:
- スペース調整機能：日本語と英語の境界に自動的にスペースを挿入
- Run境界の保持：再翻訳時の安全性を確保
- 汎用的な設計：任意のPPTXファイル、言語ペアに対応

使用方法:
    python apply_translations.py <input.pptx> <translations_dir> <output.pptx>

入力:
    translations/slide1_translations.json, slide2_translations.json, ...

対訳JSON形式:
    {
      "slide_number": 1,
      "translations": {
        "0_0_0": {
          "original": "いまさら聞けない生成",
          "translated": "A Beginner's Guide to ",
          "changed": true
        }
      }
    }
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
import json
import sys
from pathlib import Path


def get_char_type(char: str) -> str:
    """
    文字の種類を判定

    Args:
        char: 判定する文字

    Returns:
        "japanese", "english", または "other"
    """
    if '\u3040' <= char <= '\u9FFF':  # 日本語（ひらがな、カタカナ、漢字）
        return "japanese"
    if char.isalnum():  # 英数字
        return "english"
    return "other"


def insert_spaces_at_language_boundaries(text: str) -> str:
    """
    日本語と英語の境界にスペースを挿入

    日→英翻訳時にスペース不足を防ぐための処理。
    英→日再翻訳時はrun境界が保持されるため安全。

    Args:
        text: スペース調整するテキスト

    Returns:
        スペース調整後のテキスト
    """
    result = []
    prev_type = None

    for char in text:
        curr_type = get_char_type(char)

        # 日本語→英語、英語→日本語の境界にスペースを挿入
        if ((prev_type == "japanese" and curr_type == "english") or
            (prev_type == "english" and curr_type == "japanese")):
            result.append(' ')

        result.append(char)
        prev_type = curr_type

    return ''.join(result)


def needs_space_before(text: str) -> bool:
    """
    テキストの前にスペースが必要かどうかを判定

    Args:
        text: 判定するテキスト

    Returns:
        スペースが必要ならTrue
    """
    if not text:
        return False

    # 次のrunの先頭が英数字、括弧、引用符などの場合はスペースが必要
    first_char = text[0]
    return first_char.isalnum() or first_char in '([{"\'<'


def ensure_trailing_space(text: str, next_run_text: str = None) -> str:
    """
    テキストの末尾にスペースを追加（必要な場合）

    英語のrun同士が連結される時に、単語同士がくっつくのを防ぐ。
    ただし、既にスペースや句読点がある場合は追加しない。

    Args:
        text: スペースを追加するテキスト
        next_run_text: 次のrunのテキスト（オプション）

    Returns:
        スペース調整後のテキスト
    """
    if not text:
        return text

    # 既に末尾がスペースまたは句読点の場合は何もしない
    if text[-1] in ' \t\n,.;:!?")]}>':
        return text

    # 次のrunがある場合、次のrunの先頭がスペースを必要とする文字ならスペースを追加
    # ただし、次のrunの先頭が句読点の場合はスペースを追加しない
    if next_run_text and next_run_text[0] not in ' \t\n,.;:!?")]}>-':
        # 次のrunの先頭が英数字や括弧などの場合はスペースを追加
        if needs_space_before(next_run_text):
            return text + ' '

    return text


def apply_translation_to_run(run, translated: str, changed: bool, marking_color, next_run_text: str = None):
    """
    1つのテキストランに翻訳を適用

    Args:
        run: pptx.text.text.TextRunオブジェクト
        translated: 翻訳済みテキスト
        changed: 変更があるかどうか
        marking_color: マーキング用の色
        next_run_text: 次のrunのテキスト（オプション）
    """
    if changed and translated != run.text:
        # スペース調整を適用
        adjusted_text = insert_spaces_at_language_boundaries(translated)
        # run間のスペースを調整
        adjusted_text = ensure_trailing_space(adjusted_text, next_run_text)
        run.text = adjusted_text
        run.font.color.rgb = marking_color


def apply_slide_translations(slide, slide_idx: int, translations_dir: Path, marking_color):
    """
    1つのスライドに対訳を適用

    Args:
        slide: pptx.slide.Slideオブジェクト
        slide_idx: スライド番号（1から開始）
        translations_dir: 翻訳JSONファイルのあるディレクトリ
        marking_color: マーキング用の色

    Returns:
        適用した翻訳数
    """
    # 翻訳JSONファイルを読み込み
    trans_path = translations_dir / f"slide{slide_idx}_translations.json"

    if not trans_path.exists():
        print(f"Slide {slide_idx}: No translation file found, skipping...")
        return 0

    with open(trans_path, 'r', encoding='utf-8') as f:
        trans_data = json.load(f)

    applied_count = 0

    # スライドに対訳を適用
    for shape_idx, shape in enumerate(slide.shapes):
        # 通常のテキストフレームを処理
        if hasattr(shape, "text_frame"):
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                # 次のrunのテキストを取得するために、全runのリストを作成
                runs_list = list(para.runs)

                for run_idx, run in enumerate(runs_list):
                    # キーは {shape_idx}_{para_idx}_{run_idx}
                    key = f"{shape_idx}_{para_idx}_{run_idx}"

                    # 次のrunのテキストを取得（スペース調整のため）
                    next_run_text = None
                    if run_idx + 1 < len(runs_list):
                        next_key = f"{shape_idx}_{para_idx}_{run_idx + 1}"
                        if next_key in trans_data['translations']:
                            next_trans = trans_data['translations'][next_key]
                            # 翻訳されているかに関わらず、次のrunのテキストを取得
                            if next_trans['changed']:
                                next_run_text = next_trans['translated']
                            else:
                                next_run_text = next_trans['original']  # 元のテキスト（英語など）
                        else:
                            # 翻訳JSONにないrunは元のテキストを使用
                            next_run_text = runs_list[run_idx + 1].text

                    if key in trans_data['translations']:
                        trans_item = trans_data['translations'][key]

                        # 翻訳されている場合は翻訳を適用
                        if trans_item['changed']:
                            apply_translation_to_run(
                                run,
                                trans_item['translated'],
                                trans_item['changed'],
                                marking_color,
                                next_run_text
                            )
                            applied_count += 1
                        else:
                            # 翻訳されていないテキスト（元から英語など）でもスペース調整を適用
                            adjusted_text = ensure_trailing_space(
                                trans_item['original'],
                                next_run_text
                            )
                            # テキストが変更された場合のみ更新
                            if adjusted_text != run.text:
                                run.text = adjusted_text
                    else:
                        # 翻訳JSONにないrun（英語のみなど）でもスペース調整を適用
                        adjusted_text = ensure_trailing_space(
                            run.text,
                            next_run_text
                        )
                        # テキストが変更された場合のみ更新
                        if adjusted_text != run.text:
                            run.text = adjusted_text

        # テーブルを処理
        if shape.has_table:
            table = shape.table
            for row_idx in range(len(table.rows)):
                row = table.rows[row_idx]
                for col_idx in range(len(row.cells)):
                    cell = row.cells[col_idx]
                    if cell.text_frame:
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            runs_list = list(para.runs)
                            for run_idx, run in enumerate(runs_list):
                                # テーブルテキストのキー: {shape_idx}_t{row}c{col}_{para_idx}_{run_idx}
                                key = f"{shape_idx}_t{row_idx}c{col_idx}_{para_idx}_{run_idx}"

                                # 次のrunのテキストを取得
                                next_run_text = None
                                if run_idx + 1 < len(runs_list):
                                    next_key = f"{shape_idx}_t{row_idx}c{col_idx}_{para_idx}_{run_idx + 1}"
                                    if next_key in trans_data['translations']:
                                        next_trans = trans_data['translations'][next_key]
                                        next_run_text = next_trans['translated'] if next_trans['changed'] else next_trans['original']
                                    else:
                                        next_run_text = runs_list[run_idx + 1].text

                                if key in trans_data['translations']:
                                    trans_item = trans_data['translations'][key]
                                    if trans_item['changed']:
                                        apply_translation_to_run(
                                            run,
                                            trans_item['translated'],
                                            trans_item['changed'],
                                            marking_color,
                                            next_run_text
                                        )
                                        applied_count += 1

    print(f"Slide {slide_idx}: Applied {applied_count} translations")

    return applied_count


def apply_all_translations(input_pptx: str, translations_dir: str, output_pptx: str):
    """
    任意のPPTXファイルに対訳を適用

    Args:
        input_pptx: 入力PPTXファイルパス
        translations_dir: 翻訳JSONファイルのあるディレクトリ
        output_pptx: 出力PPTXファイルパス
    """
    input_path = Path(input_pptx).resolve()
    translations_path = Path(translations_dir).resolve()
    output_path = Path(output_pptx).resolve()

    if not translations_path.exists():
        print(f"Error: Translations directory not found: {translations_path}")
        sys.exit(1)

    # 出力ディレクトリを作成
    output_path.parent.mkdir(parents=True, exist_ok=True)

    print(f"Applying translations to {input_path.name}...")

    # PPTXを読み込み
    prs = Presentation(input_pptx)

    # マーキング色（くすんだ紫: RGB 128, 0, 128）
    marking_color = RGBColor(128, 0, 128)

    total_applied = 0

    # 各スライドに対訳を適用
    for slide_idx, slide in enumerate(prs.slides, start=1):
        applied = apply_slide_translations(
            slide,
            slide_idx,
            translations_path,
            marking_color
        )
        total_applied += applied

    # 保存
    prs.save(output_path)

    print(f"\n✓ Successfully applied {total_applied} translations")
    print(f"Output file: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 4:
        print("Usage: python apply_translations.py <input.pptx> <translations_dir> <output.pptx>")
        print("\nExample:")
        print("  python apply_translations.py sample.pptx translations/ converted.pptx")
        sys.exit(1)

    input_pptx = sys.argv[1]
    translations_dir = sys.argv[2]
    output_pptx = sys.argv[3]

    if not Path(input_pptx).exists():
        print(f"Error: File not found: {input_pptx}")
        sys.exit(1)

    apply_all_translations(input_pptx, translations_dir, output_pptx)
